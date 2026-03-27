#!/usr/bin/env python3
"""
Конвертація PDF (наприклад експорт Canva) у PPTX: кожна сторінка PDF — слайд з зображенням.
Зберігає візуальне відповідність оригіналу (текст з PDF як вектор/шрифт залишається растровим).

Опційно: --template шлях.pptx — беруться slide_width / slide_height з шаблону (фон макета не копіюється).
"""
from __future__ import annotations

import argparse
import io
import sys
from pathlib import Path

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches


def _blank_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        name = (layout.name or "").lower()
        if "blank" in name or "порожн" in name:
            return layout
    return prs.slide_layouts[-1]


def pdf_to_pptx(
    pdf_path: Path,
    out_path: Path,
    template_path: Path | None = None,
    zoom: float = 2.0,
) -> int:
    pdf_path = pdf_path.expanduser().resolve()
    if not pdf_path.is_file():
        print(f"PDF не знайдено: {pdf_path}", file=sys.stderr)
        return 1

    out_path = out_path.expanduser().resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)

    if template_path and template_path.is_file():
        tpl = Presentation(str(template_path))
        prs = Presentation()
        prs.slide_width = tpl.slide_width
        prs.slide_height = tpl.slide_height
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    sw = int(prs.slide_width)
    sh = int(prs.slide_height)
    layout = _blank_layout(prs)
    mat = fitz.Matrix(zoom, zoom)

    doc = fitz.open(pdf_path)
    try:
        for i in range(len(doc)):
            page = doc.load_page(i)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img_bytes = pix.tobytes("png")
            pix_w, pix_h = pix.width, pix.height

            slide = prs.slides.add_slide(layout)
            stream = io.BytesIO(img_bytes)

            pic_w = sw
            pic_h = int(pix_h * sw / pix_w) if pix_w else sh
            if pic_h > sh:
                pic_h = sh
                pic_w = int(pix_w * sh / pix_h) if pix_h else sw
            left = (sw - pic_w) // 2
            top = (sh - pic_h) // 2
            slide.shapes.add_picture(stream, left, top, width=pic_w)
    finally:
        doc.close()

    prs.save(str(out_path))
    n = len(prs.slides)
    # Windows cp1252 consoles often break on Ukrainian; keep message ASCII-safe
    print(f"Wrote {n} slides -> {out_path}")
    return 0


def main() -> int:
    p = argparse.ArgumentParser(description="PDF → PPTX (сторінки як зображення на слайдах)")
    p.add_argument("pdf", type=Path, help="Шлях до .pdf")
    p.add_argument(
        "-o",
        "--output",
        type=Path,
        default=None,
        help="Вихідний .pptx (за замовчуванням: поруч з PDF, суфікс -from-pdf.pptx)",
    )
    p.add_argument(
        "-t",
        "--template",
        type=Path,
        default=None,
        help="Опційний .pptx для розмірів слайда (ширина/висота)",
    )
    p.add_argument(
        "--zoom",
        type=float,
        default=2.0,
        help="Масштаб рендеру PDF (більше = чіткіше, важчий файл)",
    )
    args = p.parse_args()
    out = args.output
    if out is None:
        out = args.pdf.with_name(args.pdf.stem + "-from-pdf.pptx")
    return pdf_to_pptx(args.pdf, out, args.template, args.zoom)


if __name__ == "__main__":
    raise SystemExit(main())
